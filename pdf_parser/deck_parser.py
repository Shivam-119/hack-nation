from __future__ import annotations

from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


class DeckParseError(Exception):
    pass


def extract_deck_text(path: str | Path) -> str:
    path = Path(path)
    if not path.exists():
        raise DeckParseError(f"Deck file not found: {path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    raise DeckParseError(f"Unsupported deck type: {suffix} (expected .pdf or .pptx)")


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    slides: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        slides.append(f"--- Slide {i} ---\n{text}")
    return "\n\n".join(slides)


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes] {notes}")
        slides.append(f"--- Slide {i} ---\n" + "\n".join(parts))
    return "\n\n".join(slides)
